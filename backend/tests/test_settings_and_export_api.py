import json
import asyncio
import io
from types import SimpleNamespace

from fastapi.testclient import TestClient
from pptx import Presentation as PptxPresentation

from app.api.v1 import export as export_api
from app.main import app


def _sample_export_payload() -> dict:
    return {
        "presentation": {
            "presentationId": "pres-1",
            "title": "测试文稿",
            "slides": [
                {
                    "slideId": "slide-1",
                    "layoutType": "intro-slide",
                    "layoutId": "intro-slide",
                    "contentData": {"title": "封面"},
                    "components": [],
                }
            ],
        }
    }


def test_settings_api_persists_enable_vision_verification(tmp_path, monkeypatch):
    from app.core import settings_store
    from app.core.config import settings

    original_value = settings.enable_vision_verification
    original_primary_strategy = settings.content_type_primary_strategy
    original_shadow_enabled = settings.content_type_shadow_enabled
    original_confidence_threshold = settings.content_type_confidence_threshold
    settings_path = tmp_path / "settings.json"
    monkeypatch.setattr(settings_store, "_SETTINGS_FILE", settings_path)

    client = TestClient(app)

    try:
        resp = client.put(
            "/api/v1/settings",
            json={
                "enable_vision_verification": False,
                "content_type_primary_strategy": "semantic",
                "content_type_shadow_enabled": False,
                "content_type_confidence_threshold": 0.73,
            },
        )
        assert resp.status_code == 200
        assert resp.json()["enable_vision_verification"] is False
        assert resp.json()["content_type_primary_strategy"] == "semantic"
        assert resp.json()["content_type_shadow_enabled"] is False
        assert resp.json()["content_type_confidence_threshold"] == 0.73

        assert settings_path.exists()
        persisted = json.loads(settings_path.read_text(encoding="utf-8"))
        assert persisted["enable_vision_verification"] is False
        assert persisted["content_type_primary_strategy"] == "semantic"
        assert persisted["content_type_shadow_enabled"] is False
        assert persisted["content_type_confidence_threshold"] == 0.73

        get_resp = client.get("/api/v1/settings")
        assert get_resp.status_code == 200
        assert get_resp.json()["enable_vision_verification"] is False
        assert get_resp.json()["content_type_primary_strategy"] == "semantic"
        assert get_resp.json()["content_type_shadow_enabled"] is False
        assert get_resp.json()["content_type_confidence_threshold"] == 0.73
    finally:
        settings.enable_vision_verification = original_value
        settings.content_type_primary_strategy = original_primary_strategy
        settings.content_type_shadow_enabled = original_shadow_enabled
        settings.content_type_confidence_threshold = original_confidence_threshold


def test_export_pdf_returns_503_with_manual_install_hint(monkeypatch):
    from app.services.export import pdf_exporter

    async def fake_export_pdf(html_content: str):  # noqa: ARG001
        raise RuntimeError(
            "Playwright Chromium 自动安装失败。安装超时（>90s）。"
            "请手动执行: uv run playwright install chromium"
        )

    monkeypatch.setattr(pdf_exporter, "export_pdf", fake_export_pdf)

    client = TestClient(app)
    resp = client.post("/api/v1/export/pdf", json=_sample_export_payload())

    assert resp.status_code == 503
    assert "uv run playwright install chromium" in resp.json()["detail"]


def test_export_content_disposition_uses_rfc5987_for_non_ascii_title():
    header = export_api._build_content_disposition("测试汇报", "pptx")

    assert 'filename="presentation.pptx"' in header
    assert "filename*=UTF-8''%E6%B5%8B%E8%AF%95%E6%B1%87%E6%8A%A5.pptx" in header
    header.encode("latin-1")


def test_export_pdf_returns_non_ascii_filename_header(monkeypatch):
    from app.services.export import pdf_exporter

    async def fake_export_pdf(html_content: str):  # noqa: ARG001
        return b"%PDF-1.4"

    monkeypatch.setattr(pdf_exporter, "export_pdf", fake_export_pdf)

    client = TestClient(app)
    resp = client.post("/api/v1/export/pdf", json=_sample_export_payload())

    assert resp.status_code == 200
    assert 'filename="presentation.pdf"' in resp.headers["content-disposition"]
    assert "filename*=UTF-8''" in resp.headers["content-disposition"]


def test_export_pptx_sets_structured_mode_header(monkeypatch):
    from app.services.export import pptx_exporter

    monkeypatch.setattr(pptx_exporter, "export_pptx", lambda presentation: b"PPTX")

    client = TestClient(app)
    resp = client.post("/api/v1/export/pptx", json=_sample_export_payload())

    assert resp.status_code == 200
    assert resp.content == b"PPTX"
    assert resp.headers["x-zhiyan-export-mode"] == "structured"


def test_export_pptx_falls_back_to_image_mode(monkeypatch):
    from app.services.export import pptx_exporter, pptx_image_exporter

    def fake_export_pptx(presentation):  # noqa: ARG001
        raise ValueError("structured failed")

    async def fake_export_pptx_as_images(presentation):  # noqa: ARG001
        return b"FALLBACK"

    monkeypatch.setattr(pptx_exporter, "export_pptx", fake_export_pptx)
    monkeypatch.setattr(pptx_image_exporter, "export_pptx_as_images", fake_export_pptx_as_images)

    client = TestClient(app)
    resp = client.post("/api/v1/export/pptx", json=_sample_export_payload())

    assert resp.status_code == 200
    assert resp.content == b"FALLBACK"
    assert resp.headers["x-zhiyan-export-mode"] == "fallback-image"


def test_export_pptx_returns_manual_install_hint_when_fallback_unavailable(monkeypatch):
    from app.services.export import pptx_exporter, pptx_image_exporter

    def fake_export_pptx(presentation):  # noqa: ARG001
        raise ValueError("structured failed")

    async def fake_export_pptx_as_images(presentation):  # noqa: ARG001
        raise RuntimeError(
            "Playwright Chromium 自动安装失败。安装超时（>90s）。"
            "请手动执行: uv run playwright install chromium"
        )

    monkeypatch.setattr(pptx_exporter, "export_pptx", fake_export_pptx)
    monkeypatch.setattr(pptx_image_exporter, "export_pptx_as_images", fake_export_pptx_as_images)

    client = TestClient(app)
    resp = client.post("/api/v1/export/pptx", json=_sample_export_payload())

    assert resp.status_code == 503
    assert "uv run playwright install chromium" in resp.json()["detail"]


def test_export_pptx_as_images_creates_one_slide_per_screenshot(monkeypatch):
    from app.models.slide import Presentation
    from app.services.export import slide_screenshot
    from app.services.export.pptx_image_exporter import export_pptx_as_images

    screenshot_bytes = (
        b"\x89PNG\r\n\x1a\n"
        b"\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01\x08\x06\x00\x00\x00\x1f\x15\xc4\x89"
        b"\x00\x00\x00\rIDATx\x9cc```\xf8\x0f\x00\x01\x04\x01\x00]\xc2\x02~"
        b"\x00\x00\x00\x00IEND\xaeB`\x82"
    )

    async def fake_capture_slide_screenshots(presentation_dict, job_id=None):  # noqa: ARG001
        return [
            slide_screenshot.SlideScreenshot(slide_id="slide-1", png_bytes=screenshot_bytes),
            slide_screenshot.SlideScreenshot(slide_id="slide-2", png_bytes=screenshot_bytes),
        ]

    monkeypatch.setattr(slide_screenshot, "capture_slide_screenshots", fake_capture_slide_screenshots)

    payload = Presentation.model_validate(_sample_export_payload()["presentation"])
    pptx_bytes = asyncio.run(export_pptx_as_images(payload))
    prs = PptxPresentation(io.BytesIO(pptx_bytes))

    assert len(prs.slides) == 2
    for slide in prs.slides:
        assert len(slide.shapes) == 1


def test_export_pptx_as_images_rejects_empty_screenshot_list(monkeypatch):
    from app.models.slide import Presentation
    from app.services.export import slide_screenshot
    from app.services.export.pptx_image_exporter import export_pptx_as_images

    async def fake_capture_slide_screenshots(presentation_dict, job_id=None):  # noqa: ARG001
        return []

    monkeypatch.setattr(slide_screenshot, "capture_slide_screenshots", fake_capture_slide_screenshots)

    payload = Presentation.model_validate(_sample_export_payload()["presentation"])

    try:
        asyncio.run(export_pptx_as_images(payload))
    except RuntimeError as exc:
        assert str(exc) == "No slides available for fallback export"
    else:
        raise AssertionError("expected fallback exporter to reject an empty screenshot list")


def test_select_blank_layout_prefers_name_then_falls_back_to_index():
    from app.services.export.pptx_image_exporter import _select_blank_layout

    class FakeLayouts:
        def __init__(self, layouts):
            self._layouts = layouts

        def __iter__(self):
            return iter(self._layouts)

        def __getitem__(self, index):
            return self._layouts[index]

    named_blank = SimpleNamespace(name="Blank")
    index_blank = SimpleNamespace(name="Layout 6")

    prs_with_named_blank = SimpleNamespace(
        slide_layouts=FakeLayouts(
            [SimpleNamespace(name=f"Layout {idx}") for idx in range(6)] + [named_blank]
        )
    )
    assert _select_blank_layout(prs_with_named_blank) is named_blank

    layouts_without_named_blank = [SimpleNamespace(name=f"Layout {idx}") for idx in range(7)]
    layouts_without_named_blank[6] = index_blank
    prs_without_named_blank = SimpleNamespace(slide_layouts=FakeLayouts(layouts_without_named_blank))
    assert _select_blank_layout(prs_without_named_blank) is index_blank


def test_cors_preflight_allows_localhost_and_loopback():
    client = TestClient(app)

    for origin in ("http://127.0.0.1:3000", "http://localhost:5173"):
        resp = client.options(
            "/api/v1/settings",
            headers={
                "Origin": origin,
                "Access-Control-Request-Method": "GET",
                "Access-Control-Request-Headers": "content-type,x-workspace-id",
            },
        )
        assert resp.status_code == 200
        assert resp.headers.get("access-control-allow-origin") == origin


def test_cors_preflight_rejects_unknown_origin():
    client = TestClient(app)
    resp = client.options(
        "/api/v1/settings",
        headers={
            "Origin": "https://evil.example.com",
            "Access-Control-Request-Method": "GET",
        },
    )
    assert resp.status_code == 400
