"""PPTX 导出 — 从 Slide JSON 生成 .pptx 文件

使用 python-pptx 将 Slide JSON 重建为 PowerPoint 文件。
支持两种格式：
- 新版 layoutId + contentData：按结构化数据渲染
- 旧版 components：按组件坐标渲染（向后兼容）
"""

import io
import re
from typing import Any

from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

from app.models.slide import Presentation, Component
from app.services.presentations.normalizer import (
    normalize_metrics_slide_data,
    normalize_outline_slide_data,
    split_outline_sections,
)
from app.services.export.layout_rules import (
    get_bullet_with_icons_columns,
    is_bullet_icons_only_compact,
)

# 16:9 宽屏尺寸
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def _pct_to_emu(pct: float, total: int) -> int:
    """百分比转 EMU"""
    return int(total * pct / 100)


def _parse_color(color_str: str | None) -> RGBColor | None:
    """解析 CSS 颜色字符串为 RGBColor"""
    if not color_str:
        return None
    color_str = color_str.strip().lstrip("#")
    if len(color_str) == 6:
        try:
            return RGBColor(
                int(color_str[0:2], 16),
                int(color_str[2:4], 16),
                int(color_str[4:6], 16),
            )
        except ValueError:
            return None
    return None


def _get_alignment(text_align: str | None) -> int | None:
    mapping = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }
    return mapping.get(text_align)


def _is_bullet_line(line: str) -> tuple[bool, int, str]:
    """检测列表行。返回 (is_bullet, level, clean_text)"""
    nested = re.match(r'^(\s{2,}|\t+)([•\-*]|\d+[.)])\s+(.*)', line)
    if nested:
        return True, 1, nested.group(3)
    unordered = re.match(r'^[•\-*]\s+(.*)', line)
    if unordered:
        return True, 0, unordered.group(1)
    ordered = re.match(r'^\d+[.)]\s+(.*)', line)
    if ordered:
        return True, 0, ordered.group(1)
    return False, 0, line


def _add_textbox(slide_obj, left: int, top: int, width: int, height: int,
                 text: str, font_size: float = 18, bold: bool = False,
                 color: RGBColor | None = None, alignment: int | None = None):
    """添加文本框的便捷函数"""
    txBox = slide_obj.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if font_size:
        p.font.size = Pt(font_size)
    if bold:
        p.font.bold = True
    if color:
        p.font.color.rgb = color
    if alignment is not None:
        p.alignment = alignment
    return txBox


def _add_text_component(slide_obj, comp: Component) -> None:
    """添加文本组件到幻灯片（旧版）"""
    pos = comp.position
    left = _pct_to_emu(pos.x, SLIDE_WIDTH)
    top = _pct_to_emu(pos.y, SLIDE_HEIGHT)
    width = _pct_to_emu(pos.width, SLIDE_WIDTH)
    height = _pct_to_emu(pos.height, SLIDE_HEIGHT)

    txBox = slide_obj.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    content = comp.content or ""
    style = comp.style
    lines = content.split("\n")

    for i, line in enumerate(lines):
        line_stripped = line.strip()
        if not line_stripped:
            continue

        is_bullet, level, clean_text = _is_bullet_line(line_stripped)

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = clean_text

        if is_bullet:
            p.level = level
            pPr = p._pPr
            if pPr is None:
                pPr = p._p.get_or_add_pPr()
            buChar = pPr.makeelement(qn("a:buChar"), {"char": "•"})
            for old in pPr.findall(qn("a:buChar")):
                pPr.remove(old)
            for old in pPr.findall(qn("a:buNone")):
                pPr.remove(old)
            pPr.append(buChar)

        if style:
            font_size = style.font_size
            if font_size:
                p.font.size = Pt(font_size)
            if style.font_weight == "bold":
                p.font.bold = True
            color = _parse_color(style.color)
            if color:
                p.font.color.rgb = color
            alignment = _get_alignment(
                style.text_align.value if style.text_align else None
            )
            if alignment is not None:
                p.alignment = alignment


def _set_slide_background(slide_obj, bg_color: str | None) -> None:
    """设置幻灯片背景色"""
    color = _parse_color(bg_color)
    if not color:
        return
    background = slide_obj.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_placeholder_component(slide_obj, comp: Component) -> None:
    """添加占位组件（图片/图表）— 旧版"""
    pos = comp.position
    left = _pct_to_emu(pos.x, SLIDE_WIDTH)
    top = _pct_to_emu(pos.y, SLIDE_HEIGHT)
    width = _pct_to_emu(pos.width, SLIDE_WIDTH)
    height = _pct_to_emu(pos.height, SLIDE_HEIGHT)

    txBox = slide_obj.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"[{comp.type.value}: {comp.content or '占位'}]"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.CENTER

    shape = txBox
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF3, 0xF4, 0xF6)


# ---------- 新版 contentData 渲染 ----------

PRIMARY_COLOR = RGBColor(0x3B, 0x82, 0xF6)
GRAY_600 = RGBColor(0x4B, 0x55, 0x63)
GRAY_400 = RGBColor(0x9C, 0xA3, 0xAF)


def _as_text(value: Any, default: str = "") -> str:
    if isinstance(value, str):
        stripped = value.strip()
        if stripped:
            return stripped
    if isinstance(value, (int, float, bool)):
        return str(value)
    return default


def _item_text(item: Any) -> str:
    if isinstance(item, str):
        return item.strip()
    if isinstance(item, dict):
        return (
            _as_text(item.get("text"))
            or _as_text(item.get("title"))
            or _as_text(item.get("label"))
            or _as_text(item.get("challenge"))
            or _as_text(item.get("outcome"))
        )
    return ""


def _item_description(item: Any) -> str:
    if isinstance(item, dict):
        return _as_text(item.get("description"))
    return ""


def _item_icon_query(item: Any) -> str:
    if isinstance(item, dict):
        icon = item.get("icon")
        if isinstance(icon, dict):
            query = _as_text(icon.get("query"))
            if query:
                return query
    return _item_text(item)


def _icon_token(query: str) -> str:
    cleaned = "".join(ch for ch in query if ch.isalnum())
    if not cleaned:
        return "IC"
    return cleaned[:2].upper()


def _table_headers_and_rows(data: dict[str, Any]) -> tuple[list[str], list[list[str]]]:
    headers_raw = data.get("headers")
    if not isinstance(headers_raw, list):
        headers_raw = data.get("columns")
    headers = [_as_text(header) for header in headers_raw if _as_text(header)] if isinstance(headers_raw, list) else []

    rows_raw = data.get("rows")
    rows: list[list[str]] = []
    if isinstance(rows_raw, list):
        for row in rows_raw:
            if isinstance(row, list):
                values = [_as_text(cell) for cell in row]
                if len(values) < len(headers):
                    values.extend([""] * (len(headers) - len(values)))
                rows.append(values[: len(headers)])
            elif isinstance(row, dict):
                rows.append([_as_text(row.get(header)) for header in headers])
    return headers, rows


def _extract_compare_column(raw: Any, fallback: str) -> tuple[str, list[str]]:
    if not isinstance(raw, dict):
        return fallback, []
    heading = _as_text(raw.get("heading")) or _as_text(raw.get("title")) or fallback
    items_raw = raw.get("items")
    if not isinstance(items_raw, list):
        return heading, []
    items = [_item_text(item) for item in items_raw]
    return heading, [item for item in items if item]


def _extract_challenge_outcome_pairs(data: dict[str, Any]) -> list[tuple[str, str]]:
    pairs: list[tuple[str, str]] = []
    raw_items = data.get("items")
    if isinstance(raw_items, list):
        for row in raw_items:
            if isinstance(row, dict):
                pairs.append(
                    (
                        _as_text(row.get("challenge"), "内容生成中"),
                        _as_text(row.get("outcome"), "待补充"),
                    )
                )
            elif isinstance(row, str) and row.strip():
                pairs.append((row.strip(), "待补充"))
    if pairs:
        return pairs

    for side_key in ("challenge", "outcome"):
        if not isinstance(data.get(side_key), dict):
            return []
    challenge_items_raw = data["challenge"].get("items")
    outcome_items_raw = data["outcome"].get("items")
    challenge_items = [_item_text(item) for item in challenge_items_raw] if isinstance(challenge_items_raw, list) else []
    outcome_items = [_item_text(item) for item in outcome_items_raw] if isinstance(outcome_items_raw, list) else []
    count = max(len(challenge_items), len(outcome_items))
    if count == 0:
        return []
    for idx in range(count):
        pairs.append(
            (
                challenge_items[idx] if idx < len(challenge_items) and challenge_items[idx] else "内容生成中",
                outcome_items[idx] if idx < len(outcome_items) and outcome_items[idx] else "待补充",
            )
        )
    return pairs


GRAY_200 = RGBColor(0xDB, 0xE2, 0xEA)
GRAY_900 = RGBColor(0x0F, 0x17, 0x2A)

OUTLINE_ACCENT_LEFT = Inches(0.82)
OUTLINE_ACCENT_TOP = Inches(0.62)
OUTLINE_ACCENT_WIDTH = Inches(0.72)
OUTLINE_ACCENT_HEIGHT = Inches(0.06)
OUTLINE_TITLE_TOP = Inches(0.90)
OUTLINE_TITLE_WIDTH = Inches(5.4)
OUTLINE_TITLE_HEIGHT = Inches(0.70)
OUTLINE_SUBTITLE_TOP = Inches(1.62)
OUTLINE_SUBTITLE_WIDTH = Inches(5.2)
OUTLINE_SUBTITLE_HEIGHT = Inches(0.58)
OUTLINE_DIVIDER_LEFT = Inches(6.25)
OUTLINE_DIVIDER_TOP = Inches(1.18)
OUTLINE_DIVIDER_WIDTH = Inches(6.0)
OUTLINE_DIVIDER_HEIGHT = Inches(0.02)
OUTLINE_LEFT_COLUMN_LEFT = Inches(0.82)
OUTLINE_RIGHT_COLUMN_LEFT = Inches(6.90)
OUTLINE_CONTENT_TOP = Inches(2.18)
OUTLINE_COLUMN_WIDTH = Inches(5.45)
OUTLINE_COLUMN_HEIGHT = Inches(4.75)


def _add_rule(slide_obj, left: int, top: int, width: int, height: int, color: RGBColor):
    shape = slide_obj.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _render_outline_column(
    slide_obj,
    sections: list[dict[str, str]],
    start_index: int,
    left: int,
    top: int,
    width: int,
    height: int,
    accent_color: RGBColor,
) -> None:
    if not sections:
        return

    row_height = height // len(sections)
    number_width = Inches(0.55)
    gap_width = Inches(0.18)
    title_left = left + number_width + gap_width
    title_width = width - number_width - gap_width

    for offset, section in enumerate(sections):
        item_top = top + row_height * offset
        _add_rule(slide_obj, left, item_top, width, Inches(0.02), GRAY_200)
        _add_textbox(
            slide_obj,
            left,
            item_top + Inches(0.10),
            number_width,
            Inches(0.25),
            f"{start_index + offset + 1:02d}",
            font_size=11,
            bold=True,
            color=accent_color,
        )
        _add_textbox(
            slide_obj,
            title_left,
            item_top + Inches(0.05),
            title_width,
            Inches(0.45),
            section.get("title", f"Section {start_index + offset + 1}"),
            font_size=21,
            bold=True,
            color=GRAY_900,
        )
        description = section.get("description", "")
        if description:
            _add_textbox(
                slide_obj,
                title_left,
                item_top + Inches(0.55),
                title_width,
                max(row_height - Inches(0.7), Inches(0.3)),
                description,
                font_size=12,
                color=GRAY_600,
            )

def _render_content_data(slide_obj, layout_id: str, data: dict, theme_color: RGBColor | None = None) -> None:
    """根据 layout_id 和 contentData 渲染幻灯片内容"""
    color = theme_color or PRIMARY_COLOR
    d = data

    if layout_id == "intro-slide":
        _add_textbox(slide_obj,
                     Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
                     d.get("title", ""), font_size=48, bold=True, color=color,
                     alignment=PP_ALIGN.CENTER)
        if d.get("subtitle"):
            _add_textbox(slide_obj,
                         Inches(2), Inches(3.8), Inches(9), Inches(0.8),
                         d["subtitle"], font_size=24, color=GRAY_600,
                         alignment=PP_ALIGN.CENTER)
        info_parts = []
        author = _as_text(d.get("author")) or _as_text(d.get("presenter"))
        if author:
            info_parts.append(author)
        if d.get("date"):
            info_parts.append(d["date"])
        if info_parts:
            _add_textbox(slide_obj,
                         Inches(2), Inches(5.0), Inches(9), Inches(0.6),
                         " | ".join(info_parts), font_size=16, color=GRAY_400,
                         alignment=PP_ALIGN.CENTER)

    elif layout_id == "section-header":
        if d.get("section_number"):
            _add_textbox(slide_obj,
                         Inches(2), Inches(2.2), Inches(9), Inches(0.6),
                         str(d["section_number"]), font_size=18, color=color,
                         alignment=PP_ALIGN.CENTER)
        _add_textbox(slide_obj,
                     Inches(1.5), Inches(2.8), Inches(10), Inches(1.2),
                     d.get("title", ""), font_size=44, bold=True,
                     alignment=PP_ALIGN.CENTER)
        if d.get("subtitle"):
            _add_textbox(slide_obj,
                         Inches(2), Inches(4.2), Inches(9), Inches(0.8),
                         d["subtitle"], font_size=22, color=GRAY_600,
                         alignment=PP_ALIGN.CENTER)

    elif layout_id == "outline-slide":
        outline = normalize_outline_slide_data(d)
        left_sections, right_sections = split_outline_sections(outline["sections"])
        _add_rule(slide_obj, OUTLINE_ACCENT_LEFT, OUTLINE_ACCENT_TOP, OUTLINE_ACCENT_WIDTH, OUTLINE_ACCENT_HEIGHT, color)
        _add_textbox(slide_obj,
                     OUTLINE_ACCENT_LEFT, OUTLINE_TITLE_TOP, OUTLINE_TITLE_WIDTH, OUTLINE_TITLE_HEIGHT,
                     _as_text(outline.get("title"), "Outline"), font_size=32, bold=True, color=GRAY_900)
        subtitle = _as_text(outline.get("subtitle"))
        if subtitle:
            _add_textbox(slide_obj,
                         OUTLINE_ACCENT_LEFT, OUTLINE_SUBTITLE_TOP, OUTLINE_SUBTITLE_WIDTH, OUTLINE_SUBTITLE_HEIGHT,
                         subtitle, font_size=14, color=GRAY_600)
        _add_rule(slide_obj, OUTLINE_DIVIDER_LEFT, OUTLINE_DIVIDER_TOP, OUTLINE_DIVIDER_WIDTH, OUTLINE_DIVIDER_HEIGHT, GRAY_200)
        _render_outline_column(
            slide_obj,
            left_sections,
            0,
            OUTLINE_LEFT_COLUMN_LEFT,
            OUTLINE_CONTENT_TOP,
            OUTLINE_COLUMN_WIDTH,
            OUTLINE_COLUMN_HEIGHT,
            color,
        )
        _render_outline_column(
            slide_obj,
            right_sections,
            len(left_sections),
            OUTLINE_RIGHT_COLUMN_LEFT,
            OUTLINE_CONTENT_TOP,
            OUTLINE_COLUMN_WIDTH,
            OUTLINE_COLUMN_HEIGHT,
            color,
        )
    elif layout_id == "bullet-with-icons":
        _add_textbox(slide_obj,
                     Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                     d.get("title", ""), font_size=36, bold=True)
        items_source = d.get("items")
        if not isinstance(items_source, list):
            items_source = d.get("features", [])
        items = items_source if isinstance(items_source, list) else []
        columns = get_bullet_with_icons_columns(len(items))
        compact = columns == 4
        gutter = 0.18 if compact else 0.28
        content_width = 11.4
        column_width = (content_width - gutter * (columns - 1)) / columns
        base_left = 0.8

        for idx, item in enumerate(items[:4]):
            left = base_left + idx * (column_width + gutter)
            title = _item_text(item)
            desc = _item_description(item)
            icon_token = _icon_token(_item_icon_query(item))

            rule = slide_obj.shapes.add_shape(
                1,
                Inches(left),
                Inches(2.55),
                Inches(0.015),
                Inches(2.25 if compact else 2.45),
            )
            rule.fill.solid()
            rule.fill.fore_color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
            rule.line.fill.background()

            title_bg = slide_obj.shapes.add_shape(
                1,
                Inches(left + 0.18),
                Inches(2.0),
                Inches(max(column_width - 0.28, 0.8)),
                Inches(0.34),
            )
            title_bg.fill.solid()
            title_bg.fill.fore_color.rgb = RGBColor(0xEF, 0xF5, 0xFE)
            title_bg.line.fill.background()

            icon_box = slide_obj.shapes.add_shape(
                1,
                Inches(left + 0.18),
                Inches(1.22),
                Inches(0.42),
                Inches(0.42),
            )
            icon_box.fill.solid()
            icon_box.fill.fore_color.rgb = RGBColor(0xEF, 0xF6, 0xFF)
            icon_box.line.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)

            _add_textbox(
                slide_obj,
                Inches(left + 0.19),
                Inches(1.31),
                Inches(0.40),
                Inches(0.18),
                icon_token,
                font_size=10,
                bold=True,
                color=PRIMARY_COLOR,
                alignment=PP_ALIGN.CENTER,
            )

            _add_textbox(
                slide_obj,
                Inches(left + 0.18),
                Inches(1.9),
                Inches(max(column_width - 0.24, 0.8)),
                Inches(0.7),
                title,
                font_size=19 if compact else 21,
                bold=True,
                color=PRIMARY_COLOR,
            )
            if desc:
                _add_textbox(
                    slide_obj,
                    Inches(left + 0.18),
                    Inches(2.72),
                    Inches(max(column_width - 0.24, 0.8)),
                    Inches(1.15 if compact else 1.35),
                    desc,
                    font_size=11.5 if compact else 12.5,
                    color=GRAY_600,
                )
            _add_textbox(
                slide_obj,
                Inches(left + 0.18),
                Inches(5.42),
                Inches(max(column_width - 0.24, 0.8)),
                Inches(0.65),
                str(idx + 1).zfill(2),
                font_size=40 if compact else 46,
                color=RGBColor(0x11, 0x18, 0x27),
            )

    elif layout_id == "bullet-icons-only":
        _add_textbox(
            slide_obj,
            Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
            d.get("title", ""), font_size=36, bold=True
        )
        items_source = d.get("items")
        if not isinstance(items_source, list):
            items_source = d.get("features", [])
        items = items_source if isinstance(items_source, list) else []
        compact = is_bullet_icons_only_compact(len(items))
        if items:
            card_width = 5.45
            card_height = 0.98 if compact else 1.08
            start_y = 1.85
            gap_y = 0.22 if compact else 0.26
            for idx, item in enumerate(items[:8]):
                col_idx = idx % 2
                row_idx = idx // 2
                left = 0.8 + col_idx * 5.75
                top = start_y + row_idx * (card_height + gap_y)
                icon_token = _icon_token(_item_icon_query(item))

                card = slide_obj.shapes.add_shape(
                    1, Inches(left), Inches(top), Inches(card_width), Inches(card_height)
                )
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(0xF9, 0xFA, 0xFB)
                card.line.fill.background()

                accent = slide_obj.shapes.add_shape(
                    1, Inches(left + 0.28), Inches(top + 0.28), Inches(1.0), Inches(0.48)
                )
                accent.fill.solid()
                accent.fill.fore_color.rgb = RGBColor(0xDB, 0xEA, 0xFE)
                accent.line.fill.background()

                number_box = slide_obj.shapes.add_shape(
                    1, Inches(left + 0.34), Inches(top + 0.18), Inches(0.72), Inches(0.72)
                )
                number_box.fill.solid()
                number_box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                number_box.line.color.rgb = RGBColor(0xDB, 0xEA, 0xFE)

                _add_textbox(
                    slide_obj,
                    Inches(left + 0.39), Inches(top + 0.34), Inches(0.62), Inches(0.2),
                    icon_token, font_size=12, bold=True, color=PRIMARY_COLOR,
                    alignment=PP_ALIGN.CENTER,
                )
                _add_textbox(
                    slide_obj,
                    Inches(left + 1.28), Inches(top + 0.2), Inches(card_width - 1.5), Inches(0.18),
                    str(idx + 1).zfill(2), font_size=9, bold=True, color=GRAY_600,
                )
                _add_textbox(
                    slide_obj,
                    Inches(left + 1.28), Inches(top + 0.42), Inches(card_width - 1.5), Inches(0.38),
                    _item_text(item), font_size=21 if compact else 24, bold=True,
                )

    elif layout_id == "numbered-bullets":
        _add_textbox(slide_obj,
                     Inches(0.8), Inches(0.5), Inches(11), Inches(1.0),
                     d.get("title", ""), font_size=36, bold=True)
        items = d.get("items") or d.get("steps") or []
        y = Inches(1.8)
        for i, item in enumerate(items[:6]):
            text = _item_text(item)
            _add_textbox(slide_obj,
                         Inches(0.8), y, Inches(0.6), Inches(0.5),
                         str(i + 1), font_size=24, bold=True, color=color)
            _add_textbox(slide_obj,
                         Inches(1.6), y, Inches(10), Inches(0.5),
                         text, font_size=22)
            desc = _item_description(item)
            if desc:
                _add_textbox(slide_obj,
                             Inches(1.6), y + Inches(0.5), Inches(10), Inches(0.4),
                             desc, font_size=16, color=GRAY_600)
                y += Inches(1.1)
            else:
                y += Inches(0.7)

    elif layout_id == "metrics-slide":
        normalized = normalize_metrics_slide_data(d)
        if normalized is None:
            return

        _add_textbox(slide_obj,
                     Inches(0.8), Inches(0.5), Inches(11), Inches(1.0),
                     normalized.get("title", ""), font_size=36, bold=True)
        metrics = normalized.get("metrics") or []
        conclusion = _as_text(normalized.get("conclusion"))
        conclusion_brief = _as_text(normalized.get("conclusionBrief"))
        has_executive_summary = bool(conclusion or conclusion_brief)
        count = min(len(metrics), 4)
        if has_executive_summary:
            summary_box = slide_obj.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.8), Inches(1.45), Inches(11.6), Inches(1.75),
            )
            summary_box.fill.solid()
            summary_box.fill.fore_color.rgb = RGBColor(0xF4, 0xF8, 0xFF)
            summary_box.line.color.rgb = RGBColor(0xD6, 0xE4, 0xFF)
            if conclusion:
                _add_textbox(slide_obj,
                             Inches(1.05), Inches(1.72), Inches(10.9), Inches(0.48),
                             conclusion, font_size=24, bold=True, color=GRAY_900)
            if conclusion_brief:
                _add_textbox(slide_obj,
                             Inches(1.05), Inches(2.28), Inches(10.7), Inches(0.52),
                             conclusion_brief, font_size=13, color=GRAY_600)
            if count > 0:
                card_width = 11.0 / count
                for i, m in enumerate(metrics[:4]):
                    x = Inches(0.8 + i * card_width)
                    _add_textbox(slide_obj,
                                 x + Inches(0.15), Inches(4.1), Inches(card_width - 0.3), Inches(0.7),
                                 m.get("value", ""), font_size=36, bold=True, color=color)
                    _add_textbox(slide_obj,
                                 x + Inches(0.15), Inches(4.82), Inches(card_width - 0.3), Inches(0.45),
                                 m.get("label", ""), font_size=16, bold=True)
                    if m.get("description"):
                        _add_textbox(slide_obj,
                                     x + Inches(0.15), Inches(5.28), Inches(card_width - 0.3), Inches(0.52),
                                     m["description"], font_size=12, color=GRAY_600)
        elif count > 0:
            card_width = 10.0 / count
            for i, m in enumerate(metrics[:4]):
                x = Inches(1.2 + i * card_width)
                _add_textbox(slide_obj,
                             x, Inches(2.3), Inches(card_width), Inches(0.9),
                             m.get("value", ""), font_size=42, bold=True, color=color,
                             alignment=PP_ALIGN.CENTER)
                _add_textbox(slide_obj,
                             x, Inches(3.2), Inches(card_width), Inches(0.5),
                             m.get("label", ""), font_size=18, bold=True,
                             alignment=PP_ALIGN.CENTER)
                if m.get("description"):
                    _add_textbox(slide_obj,
                                 x, Inches(3.75), Inches(card_width), Inches(0.8),
                                 m["description"], font_size=13, color=GRAY_600,
                                 alignment=PP_ALIGN.CENTER)

    elif layout_id == "metrics-with-image":
        _add_textbox(slide_obj,
                     Inches(0.8), Inches(0.8), Inches(5.5), Inches(0.9),
                     d.get("title", ""), font_size=32, bold=True)
        metrics = d.get("metrics") or []
        y = Inches(2.0)
        for m in metrics[:4]:
            _add_textbox(slide_obj,
                         Inches(0.8), y, Inches(2), Inches(0.6),
                         m.get("value", ""), font_size=32, bold=True, color=color)
            _add_textbox(slide_obj,
                         Inches(3.0), y + Inches(0.05), Inches(3.5), Inches(0.5),
                         m.get("label", ""), font_size=16, color=GRAY_600)
            y += Inches(0.8)
        # 右侧图片占位
        _add_textbox(slide_obj,
                     Inches(7.0), Inches(0.8), Inches(5.5), Inches(5.9),
                     "[图片占位]", font_size=14, color=GRAY_400,
                     alignment=PP_ALIGN.CENTER)

    elif layout_id == "chart-with-bullets":
        _add_textbox(slide_obj,
                     Inches(0.8), Inches(0.5), Inches(11), Inches(1.0),
                     d.get("title", ""), font_size=36, bold=True)
        _add_textbox(slide_obj,
                     Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0),
                     "[图表占位]", font_size=14, color=GRAY_400,
                     alignment=PP_ALIGN.CENTER)
        bullets = d.get("bullets") or []
        y = Inches(2.0)
        for b in bullets[:5]:
            text = b.get("text", "") if isinstance(b, dict) else str(b)
            _add_textbox(slide_obj,
                         Inches(7.0), y, Inches(5.5), Inches(0.5),
                         f"• {text}", font_size=18)
            y += Inches(0.6)

    elif layout_id == "table-info":
        _add_textbox(slide_obj,
                     Inches(0.8), Inches(0.5), Inches(11), Inches(1.0),
                     d.get("title", ""), font_size=36, bold=True)
        headers, rows = _table_headers_and_rows(d)
        if headers and rows:
            n_cols = len(headers)
            n_rows = min(len(rows), 8) + 1
            col_w = min(10.0 / n_cols, 3.0)
            table_shape = slide_obj.shapes.add_table(
                n_rows, n_cols,
                Inches(0.8), Inches(1.8),
                Inches(col_w * n_cols), Inches(0.5 * n_rows),
            )
            table = table_shape.table
            for ci, col_name in enumerate(headers):
                cell = table.cell(0, ci)
                cell.text = str(col_name)
                for p in cell.text_frame.paragraphs:
                    p.font.bold = True
                    p.font.size = Pt(14)
            for ri, row in enumerate(rows[:8]):
                for ci in range(n_cols):
                    cell = table.cell(ri + 1, ci)
                    cell.text = row[ci] if ci < len(row) else ""
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(13)

    elif layout_id == "two-column-compare":
        _add_textbox(slide_obj,
                     Inches(0.8), Inches(0.5), Inches(11), Inches(1.0),
                     d.get("title", ""), font_size=36, bold=True)
        left_col, right_col = _extract_compare_column(d.get("left"), "要点 A"), _extract_compare_column(d.get("right"), "要点 B")
        if not left_col[1] and not right_col[1]:
            left_col = _extract_compare_column(d.get("challenge"), "要点 A")
            right_col = _extract_compare_column(d.get("outcome"), "要点 B")
        for (col_title, col_items), x_offset in [(left_col, 0.8), (right_col, 7.0)]:
            _add_textbox(slide_obj,
                         Inches(x_offset), Inches(1.8), Inches(5.0), Inches(0.6),
                         col_title, font_size=24, bold=True, color=color)
            y = Inches(2.6)
            for text in col_items[:5]:
                _add_textbox(slide_obj,
                             Inches(x_offset), y, Inches(5.0), Inches(0.5),
                             f"• {text}", font_size=18)
                y += Inches(0.6)

    elif layout_id == "image-and-description":
        _add_textbox(slide_obj,
                     Inches(0.5), Inches(0.5), Inches(6.0), Inches(6.5),
                     "[图片占位]", font_size=14, color=GRAY_400,
                     alignment=PP_ALIGN.CENTER)
        _add_textbox(slide_obj,
                     Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.9),
                     d.get("title", ""), font_size=32, bold=True)
        if d.get("description"):
            _add_textbox(slide_obj,
                         Inches(7.0), Inches(2.8), Inches(5.5), Inches(3.5),
                         d["description"], font_size=18, color=GRAY_600)

    elif layout_id == "timeline":
        _add_textbox(slide_obj,
                     Inches(0.8), Inches(0.5), Inches(11), Inches(1.0),
                     d.get("title", ""), font_size=36, bold=True)
        events = d.get("events") or d.get("items") or []
        count = min(len(events), 5)
        if count > 0:
            step = 10.0 / count
            for i, evt in enumerate(events[:5]):
                x = Inches(1.0 + i * step)
                _add_textbox(slide_obj,
                             x, Inches(2.5), Inches(step - 0.3), Inches(0.4),
                             evt.get("date", ""), font_size=14, color=color, bold=True,
                             alignment=PP_ALIGN.CENTER)
                _add_textbox(slide_obj,
                             x, Inches(3.0), Inches(step - 0.3), Inches(0.5),
                             evt.get("title", ""), font_size=18, bold=True,
                             alignment=PP_ALIGN.CENTER)
                if evt.get("description"):
                    _add_textbox(slide_obj,
                                 x, Inches(3.6), Inches(step - 0.3), Inches(0.6),
                                 evt["description"], font_size=13, color=GRAY_600,
                                 alignment=PP_ALIGN.CENTER)

    elif layout_id == "quote-slide":
        _add_textbox(slide_obj,
                     Inches(1.5), Inches(2.0), Inches(10), Inches(2.5),
                     f'"{d.get("quote", "")}"', font_size=32, bold=False,
                     alignment=PP_ALIGN.CENTER)
        author = _as_text(d.get("author")) or _as_text(d.get("attribution"))
        context = _as_text(d.get("context"))
        footer = " · ".join(part for part in (author, context) if part)
        if footer:
            _add_textbox(slide_obj,
                         Inches(2), Inches(5.0), Inches(9), Inches(0.6),
                         f"— {footer}", font_size=18, color=GRAY_400,
                         alignment=PP_ALIGN.CENTER)

    elif layout_id == "challenge-outcome":
        _add_textbox(slide_obj,
                     Inches(0.8), Inches(0.5), Inches(11), Inches(1.0),
                     d.get("title", ""), font_size=36, bold=True)
        pairs = _extract_challenge_outcome_pairs(d)
        _add_textbox(slide_obj,
                     Inches(0.8), Inches(1.8), Inches(5.0), Inches(0.6),
                     "挑战", font_size=24, bold=True, color=RGBColor(0xDC, 0x26, 0x26))
        _add_textbox(slide_obj,
                     Inches(7.0), Inches(1.8), Inches(5.0), Inches(0.6),
                     "方案", font_size=24, bold=True, color=RGBColor(0x16, 0xA3, 0x4A))
        y = Inches(2.6)
        for challenge_text, outcome_text in pairs[:5]:
            _add_textbox(slide_obj,
                         Inches(0.8), y, Inches(5.0), Inches(0.5),
                         f"• {challenge_text}", font_size=18)
            _add_textbox(slide_obj,
                         Inches(7.0), y, Inches(5.0), Inches(0.5),
                         f"• {outcome_text}", font_size=18)
            y += Inches(0.6)

    elif layout_id == "thank-you":
        _add_textbox(slide_obj,
                     Inches(1.5), Inches(2.2), Inches(10), Inches(1.5),
                     d.get("title", "谢谢"), font_size=48, bold=True, color=color,
                     alignment=PP_ALIGN.CENTER)
        if d.get("subtitle"):
            _add_textbox(slide_obj,
                         Inches(2), Inches(4.0), Inches(9), Inches(0.8),
                         d["subtitle"], font_size=22, color=GRAY_600,
                         alignment=PP_ALIGN.CENTER)
        contact = _as_text(d.get("contact")) or _as_text(d.get("contact_info"))
        if contact:
            _add_textbox(slide_obj,
                         Inches(2), Inches(5.2), Inches(9), Inches(0.6),
                         contact, font_size=16, color=GRAY_400,
                         alignment=PP_ALIGN.CENTER)

    else:
        _add_textbox(slide_obj,
                     Inches(0.8), Inches(0.5), Inches(11), Inches(1.0),
                     d.get("title", f"[{layout_id}]"), font_size=36, bold=True)


def export_pptx(presentation: Presentation) -> bytes:
    """将 Presentation JSON 导出为 PPTX 字节"""
    prs = PptxPresentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]
    theme_color = _parse_color(
        presentation.theme.primary_color if presentation.theme else None
    )

    for slide_data in presentation.slides:
        slide_obj = prs.slides.add_slide(blank_layout)

        if presentation.theme and presentation.theme.background_color:
            _set_slide_background(slide_obj, presentation.theme.background_color)

        if slide_data.content_data and slide_data.layout_id:
            _render_content_data(slide_obj, slide_data.layout_id, slide_data.content_data, theme_color)
        else:
            for comp in slide_data.components:
                if comp.type.value == "text":
                    _add_text_component(slide_obj, comp)
                else:
                    _add_placeholder_component(slide_obj, comp)

        if slide_data.speaker_notes:
            notes_slide = slide_obj.notes_slide
            notes_slide.notes_text_frame.text = slide_data.speaker_notes

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()
