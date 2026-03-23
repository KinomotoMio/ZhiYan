"""PPTX export fallback that embeds each slide screenshot as a full-page image."""

from __future__ import annotations

import io

from pptx import Presentation as PptxPresentation
from pptx.util import Inches

from app.models.slide import Presentation


SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def _select_blank_layout(prs: PptxPresentation):
    blank = next((layout for layout in prs.slide_layouts if layout.name == "Blank"), None)
    if blank is not None:
        return blank
    return prs.slide_layouts[6]


async def export_pptx_as_images(presentation: Presentation) -> bytes:
    presentation_dict = presentation.model_dump(by_alias=True)

    from app.services.export.slide_screenshot import capture_slide_screenshots

    screenshots = await capture_slide_screenshots(presentation_dict)
    if not screenshots:
        raise RuntimeError("No slides available for fallback export")

    prs = PptxPresentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = _select_blank_layout(prs)
    for screenshot in screenshots:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            io.BytesIO(screenshot.png_bytes),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
